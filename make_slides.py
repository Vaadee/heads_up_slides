from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

import json
import random

# Load the words from the json file
with open('games_helper.json') as f:
    heads_up_words = json.load(f)["heads_up_words"]

team1_round1_words = []
team1_round2_words = []
team2_round1_words = []
team2_round2_words = []
bonus_round_team1 = []
bonus_round_team2 = []

for difficulty in heads_up_words:
    # randomly choose 10 words from the difficulty
    words = random.sample(heads_up_words[difficulty],20)
    rest_of_them = random.sample([word for word in heads_up_words[difficulty] if word not in words],10)
    assert len(rest_of_them) == 10

    # shuffle the words
    random.shuffle(words)
    random.shuffle(rest_of_them)

    # add the first 5 words to team 1 and the last 5 words to team 2
    team1_round1_words += [word.title() for word in words[:5]]
    team1_round2_words += [word.title() for word in words[5:10]]
    team2_round1_words += [word.title() for word in words[10:15]]
    team2_round2_words += [word.title() for word in words[15:20]]
    bonus_round_team1 += [word.title() for word in rest_of_them[:5]]
    bonus_round_team2 += [word.title() for word in rest_of_them[5:]]
# Create a presentation object
prs = Presentation()

# Set slide dimensions for 16:9 ratio
slide_width = Inches(16)
slide_height = Inches(9)
prs.slide_width = slide_width
prs.slide_height = slide_height

# Create a custom blank slide layout
slide_master = prs.slide_masters[0]
layout = slide_master.slide_layouts[6]  # Choosing a blank layout

# Function to add a slide with centered text
def add_centered_text_slide(text,team_num=None,font_size=84,slide_layout="words"):
    slide = prs.slides.add_slide(layout)

    if slide_layout=="title":
        left = Inches(0)
        top = Inches(0)
        width = Inches(8)
        height = Inches(9)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(33, 67, 146)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left+width, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(252, 83, 115)

    text_box_width = Inches(1.6*2)
    text_box_height = Inches(.9*2)
    text_box_left = (slide_width - text_box_width) / 2  # Centering horizontally
    text_box_top = (slide_height - text_box_height) / 2  # Centering vertically
    text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
    text_frame = text_box.text_frame
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True

    if team_num is not None:
        if team_num == 1:
            if slide_layout=="team_name":
                # underline the team name
                text_frame.paragraphs[0].font.underline = True
                text_frame.paragraphs[1].font.color.rgb = RGBColor(255,255,255)
                p.font.color.rgb = RGBColor(255,255,255)
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(33, 67, 146)
            else:
                p.font.color.rgb = RGBColor(33, 67, 146)
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(255,255,255)
        elif team_num == 2:
            if slide_layout=="team_name":
                # underline the team name
                text_frame.paragraphs[0].font.underline = True
                text_frame.paragraphs[1].font.color.rgb = RGBColor(255,255,255)
                p.font.color.rgb = RGBColor(255,255,255)
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(252, 83, 115)
            else:
                p.font.color.rgb = RGBColor(252, 83, 115)
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(255,255,255)
    else:
        p.font.color.rgb = RGBColor(255,255,255)
        text_frame.paragraphs[1].font.color.rgb = RGBColor(255,255,255)
        

# Add title slide
add_centered_text_slide("Heads Up!\nTeam Challenge",slide_layout="title")

add_centered_text_slide("Team 1\nRound 1",team_num=1,font_size=108,slide_layout="team_name")


# Add slides for Team 1 and Team 2
for word in team1_round1_words:
    add_centered_text_slide(word,team_num=1)

add_centered_text_slide("Team 2\nRound 1",team_num=2,font_size=108,slide_layout="team_name")

for word in team2_round1_words:
    add_centered_text_slide(word,team_num=2)

add_centered_text_slide("Team 1\nRound 2",team_num=1,font_size=108,slide_layout="team_name")

for word in team1_round2_words:
    add_centered_text_slide(word,team_num=1)

add_centered_text_slide("Team 2\nRound 2",team_num=2,font_size=108,slide_layout="team_name")

for word in team2_round2_words:
    add_centered_text_slide(word,team_num=2)

add_centered_text_slide("Team 1\nBonus Round",team_num=1,font_size=108,slide_layout="team_name")

for word in bonus_round_team1:
    add_centered_text_slide(word,team_num=1)

add_centered_text_slide("Team 2\nBonus Round",team_num=2,font_size=108,slide_layout="team_name")

for word in bonus_round_team2:
    add_centered_text_slide(word,team_num=2)

add_centered_text_slide("Thank You!\nHave a great day!",font_size=108,slide_layout="title")


# Save the presentation
prs.save('Heads_Up_Game_16_9.pptx')